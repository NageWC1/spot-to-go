from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x2E, 0x3D, 0x6B)   # dark navy  – headings / accents
STEEL  = RGBColor(0x4A, 0x6F, 0xA5)   # steel blue – sub-headings
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xFA)   # very light blue – slide background
DARK   = RGBColor(0x22, 0x22, 0x22)   # near-black – body text
GREEN  = RGBColor(0x27, 0xAE, 0x60)   # done / tick
ORANGE = RGBColor(0xE6, 0x7E, 0x22)   # in-progress

SCREENSHOT_DIR = os.path.join(
    os.path.dirname(__file__), "..", "screenshots app"
)

# ── helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def add_title_block(slide, title, subtitle=None):
    """Navy banner across the top with white title text."""
    banner = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(1.35)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    tf = banner.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = WHITE

    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(1.4),
                    Inches(9.3), Inches(0.45),
                    font_size=14, color=STEEL)

def bullet_block(slide, items, left, top, width, height,
                 font_size=15, title=None, title_color=NAVY):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    if title:
        p   = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(font_size + 2)
        r.font.bold  = True
        r.font.color.rgb = title_color

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"  •  {item}"
        r.font.size  = Pt(font_size)
        r.font.color.rgb = DARK

def add_image(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)

def add_note(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def status_pill(slide, label, color, left, top):
    box = slide.shapes.add_shape(1, left, top, Inches(1.6), Inches(0.35))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = label
    r.font.size  = Pt(11)
    r.font.bold  = True
    r.font.color.rgb = WHITE

# ── build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, LIGHT)

hero = s1.shapes.add_shape(1, Inches(0), Inches(2.0), Inches(10), Inches(3.5))
hero.fill.solid()
hero.fill.fore_color.rgb = NAVY
hero.line.fill.background()

add_textbox(s1, "Spot To Go",
            Inches(0.5), Inches(2.3), Inches(9), Inches(1.1),
            font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s1, "Android App — Development Progress Presentation",
            Inches(0.5), Inches(3.4), Inches(9), Inches(0.6),
            font_size=18, color=RGBColor(0xCC, 0xD9, 0xF0), align=PP_ALIGN.CENTER)

add_textbox(s1, "Academic Project  |  Kotlin + Jetpack Compose  |  Google APIs",
            Inches(0.5), Inches(4.0), Inches(9), Inches(0.5),
            font_size=13, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

add_textbox(s1, "June 2026",
            Inches(0.5), Inches(6.8), Inches(9), Inches(0.4),
            font_size=12, color=STEEL, align=PP_ALIGN.CENTER)

add_note(s1,
    "Welcome everyone. This presentation covers the end-to-end development progress "
    "of Spot To Go — an Android application that helps users discover nearby restaurants "
    "on a live map, preview them via video, and get directions. Today we'll walk through "
    "the Google Cloud API setup, Android Studio configuration including the virtual device "
    "challenge we hit, the project's file structure, and a live demo of the three working screens.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, LIGHT)
add_title_block(s2, "Project Overview",
                "What is Spot To Go and why was it built?")

bullet_block(s2, [
    "Discover nearby restaurants on an interactive Google Map",
    "Tap any restaurant marker to see its name, cuisine, rating, and distance",
    "Watch a real video preview (YouTube) directly from the detail screen",
    "Launch turn-by-turn navigation to the restaurant via Google Maps",
    "Optional login — 'Continue as Guest' lets anyone explore immediately",
], Inches(0.4), Inches(1.9), Inches(5.6), Inches(4.2),
   font_size=15, title="Core App Purpose")

bullet_block(s2, [
    "Language:  Kotlin",
    "UI Framework:  Jetpack Compose",
    "Maps:  Google Maps SDK for Android",
    "Location:  FusedLocationProviderClient",
    "Restaurant Data:  Seeded (Places API next step)",
    "IDE:  Android Studio Quail 1",
], Inches(6.3), Inches(1.9), Inches(3.4), Inches(4.0),
   font_size=13, title="Tech Stack", title_color=STEEL)

add_note(s2,
    "Spot To Go is an academic Android project. The core idea is simple: open the app, "
    "see a live map centred on your GPS position, see restaurant markers nearby, tap one "
    "to read the details, watch a preview video, or get directions. "
    "We chose Kotlin because it is the modern default in Android Studio. "
    "Jetpack Compose replaced XML layouts — it is declarative and much faster to write. "
    "The login screen exists to demonstrate navigation flow; for this demo phase 'Continue as Guest' "
    "bypasses authentication so we can focus on the map and restaurant features.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Google Cloud Console: API Setup
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, LIGHT)
add_title_block(s3, "Step 1 — Google Cloud Console: API Setup",
                "Enabling Maps SDK and Places API for the Android app")

steps = [
    "Go to console.cloud.google.com → Create new project (e.g. 'SpotToGo')",
    "Enable APIs & Services → Search and enable:  Maps SDK for Android",
    "Enable APIs & Services → Search and enable:  Places API",
    "Go to Credentials → Create API Key → restrict to Android apps",
    "Copy the API key — it will be stored in local.properties (never committed to Git)",
    "In AndroidManifest.xml: add meta-data tag referencing ${MAPS_API_KEY}",
    "In app/build.gradle.kts: read local.properties and inject via manifestPlaceholders",
]

bullet_block(s3, steps,
             Inches(0.4), Inches(1.85), Inches(9.2), Inches(4.6),
             font_size=14, title="Setup Steps")

# security note box
note_box = s3.shapes.add_shape(1, Inches(0.4), Inches(6.1), Inches(9.2), Inches(0.8))
note_box.fill.solid()
note_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
note_box.line.color.rgb = RGBColor(0xE6, 0xAC, 0x00)
tf = note_box.text_frame
tf.word_wrap = True
p  = tf.paragraphs[0]
r  = p.add_run()
r.text = ("  Security note:  local.properties is listed in .gitignore — "
          "the API key is never pushed to version control.")
r.font.size  = Pt(12)
r.font.color.rgb = RGBColor(0x7A, 0x60, 0x00)

add_note(s3,
    "The first real setup step before writing a single line of Android code was creating "
    "the project in Google Cloud Console and enabling two APIs. "
    "Maps SDK for Android renders the interactive map tile layer. "
    "Places API (Nearby Search) will return live restaurant data once we move beyond "
    "the current seeded data phase. "
    "The API key is injected at build time via manifestPlaceholders in Gradle, "
    "reading from local.properties. This file is gitignored so the key is never "
    "accidentally exposed in version control — an important security practice even for "
    "academic projects.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Android Studio Setup
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, LIGHT)
add_title_block(s4, "Step 2 — Android Studio Setup",
                "Installing the IDE and scaffolding the project")

bullet_block(s4, [
    "Installed Android Studio Quail 1 (latest stable release)",
    "Created new project:  'Empty Activity' template  →  Jetpack Compose",
    "Package name:  com.example.spottogo     Min SDK:  API 24 (Android 7.0+)",
    "Added dependencies in libs.versions.toml (version catalogue):",
    "      maps-compose 4.3.3,  play-services-maps 18.2.0",
    "      play-services-location 21.0.1,  accompanist-permissions 0.34.0",
    "      navigation-compose 2.7.7,  material-icons-extended (via Compose BOM)",
    "Compose BOM 2026.02.01 manages all Compose library versions centrally",
    "Fixed Gradle sync bug: kotlin-android plugin was missing from libs.versions.toml",
], Inches(0.4), Inches(1.85), Inches(9.2), Inches(4.8),
   font_size=14, title="Installation & Configuration Steps")

add_note(s4,
    "Android Studio Quail 1 is the current stable release. We chose the Empty Activity "
    "template with Jetpack Compose because it gives a clean starting point without "
    "legacy XML boilerplate. "
    "The version catalogue in libs.versions.toml is a Gradle best practice: "
    "all library versions are declared in one place so they can be updated "
    "without hunting through multiple build files. "
    "One issue we hit early: the kotlin-android plugin was missing from the catalogue, "
    "which silently prevented the Kotlin compiler plugin from loading during Gradle sync. "
    "Adding it resolved the sync failure. "
    "We also hit a material-icons import error because material-icons-extended was not "
    "declared as a dependency — adding it fixed all icon resolution errors in the IDE.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Virtual Device Setup (AVD + VT-x issue)
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, LIGHT)
add_title_block(s5, "Step 3 — Virtual Device (Emulator) Setup",
                "Creating an Android Virtual Device and resolving the BIOS virtualisation issue")

bullet_block(s5, [
    "Open Android Studio  →  Device Manager  →  Create Virtual Device",
    "Selected:  Pixel 6 (Medium Phone)  —  API 35 system image downloaded",
    "First run attempt: emulator failed to start — Run button stayed greyed out",
], Inches(0.4), Inches(1.85), Inches(9.2), Inches(1.3), font_size=14,
   title="Creating the AVD")

# problem box
prob = s5.shapes.add_shape(1, Inches(0.4), Inches(3.3), Inches(4.4), Inches(2.8))
prob.fill.solid()
prob.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
prob.line.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
tf = prob.text_frame; tf.word_wrap = True
p  = tf.paragraphs[0]; r = p.add_run()
r.text = "Problem Encountered"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
for line in [
    "CPU virtualisation (VT-x / AMD-V) was disabled in BIOS",
    "Confirmed via PowerShell:  HyperVRequirementVirtualizationFirmwareEnabled = False",
    "Android Emulator requires hardware virtualisation to run",
    "Gradle sync error (kotlin plugin) also kept the Run button disabled",
]:
    p2 = tf.add_paragraph(); r2 = p2.add_run()
    r2.text = f"  •  {line}"
    r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0x7B, 0x24, 0x1C)

# fix box
fix = s5.shapes.add_shape(1, Inches(5.2), Inches(3.3), Inches(4.4), Inches(2.8))
fix.fill.solid()
fix.fill.fore_color.rgb = RGBColor(0xE9, 0xF7, 0xEF)
fix.line.color.rgb = RGBColor(0x27, 0xAE, 0x60)
tf2 = fix.text_frame; tf2.word_wrap = True
p3  = tf2.paragraphs[0]; r3 = p3.add_run()
r3.text = "Fix Applied"
r3.font.size = Pt(14); r3.font.bold = True; r3.font.color.rgb = RGBColor(0x1E, 0x84, 0x49)
for line in [
    "Entered BIOS settings on restart",
    "Enabled Intel VT-x (Intel Virtualisation Technology)",
    "Saved & restarted the machine",
    "Re-ran AEHD silent installer as Administrator",
    "Emulator launched successfully — app deployed and running",
]:
    p4 = tf2.add_paragraph(); r4 = p4.add_run()
    r4.text = f"  •  {line}"
    r4.font.size = Pt(12); r4.font.color.rgb = RGBColor(0x1A, 0x5C, 0x36)

add_note(s5,
    "This was the most significant setup blocker. The Android emulator uses hardware "
    "virtualisation (Intel VT-x or AMD-V) to run efficiently. On this machine, that "
    "feature was disabled at the BIOS level. "
    "We diagnosed it by running a PowerShell command that reported "
    "HyperVRequirementVirtualizationFirmwareEnabled as False — meaning the BIOS was the "
    "only remaining blocker; Hyper-V and Windows Hypervisor Platform were already disabled. "
    "The fix was to reboot into BIOS, find the Intel Virtualisation Technology setting, "
    "enable it, save and restart. After that we re-ran the Android Emulator Hypervisor "
    "Driver (AEHD) silent installer as Administrator to register the hypervisor. "
    "The emulator then launched and the app deployed successfully.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Project Structure
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, LIGHT)
add_title_block(s6, "Project Structure",
                "How the codebase is organised — one folder per screen")

img6 = os.path.join(SCREENSHOT_DIR, "project structure and files.png")
add_image(s6, img6, Inches(0.3), Inches(1.5), Inches(4.8), Inches(5.5))

bullet_block(s6, [
    "MainActivity.kt  —  NavHost wiring all screens",
    "data/Restaurant.kt  —  data class + 5 seeded restaurants",
    "ui/login/LoginScreen.kt  —  login & guest entry",
    "ui/map/MapScreen.kt  —  map, GPS, markers, search bar",
    "ui/detail/RestaurantDetailScreen.kt  —  info, video, directions",
    "",
    "AndroidManifest.xml  —  permissions + API key injection",
    "app/build.gradle.kts  —  dependencies & manifestPlaceholders",
    "gradle/libs.versions.toml  —  version catalogue",
    "local.properties  —  API keys (gitignored)",
], Inches(5.3), Inches(1.6), Inches(4.4), Inches(5.0),
   font_size=13, title="Key Files")

add_note(s6,
    "The screenshot on the left shows the actual Android Studio file tree and editor. "
    "The project follows a simple feature-folder structure: each screen lives in its own "
    "sub-package under ui/. The data layer is a single file — Restaurant.kt holds both "
    "the data class and the repository with 5 hardcoded seed restaurants. "
    "MainActivity.kt is the entry point and navigation host; it uses Jetpack Navigation "
    "Compose to wire the three destinations: Login, Map, and Detail. "
    "The Gradle version catalogue (libs.versions.toml) keeps all dependency versions "
    "in one place, which is the modern recommended approach over hardcoding versions "
    "directly in build.gradle files.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Login Screen
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, LIGHT)
add_title_block(s7, "Screen 1 — Login",
                "Entry point and navigation gateway to the main app")

img7 = os.path.join(SCREENSHOT_DIR, "login page.png")
add_image(s7, img7, Inches(0.3), Inches(1.5), Inches(2.5), Inches(5.6))

bullet_block(s7, [
    "Displays app name 'Spot To Go' and tagline",
    "Email + Password fields for future Firebase Auth integration",
    "'Login' button — primary action (auth logic to be added in Phase 6)",
    "'Continue as Guest' — bypasses login entirely for the demo",
    "",
    "Navigation flow:",
    "  Login button  →  MapScreen",
    "  Continue as Guest  →  MapScreen",
    "",
    "Design intent:",
    "  Clean, minimal layout — no distracting elements",
    "  Steel-blue colour palette consistent across all screens",
    "  Guest option prioritised for academic demo purposes",
], Inches(3.1), Inches(1.6), Inches(6.5), Inches(5.4),
   font_size=14, title="What This Screen Does")

add_note(s7,
    "The login screen is deliberately minimal. Both the Login button and Continue as Guest "
    "navigate to the same MapScreen at this stage — the distinction is that Login will "
    "eventually validate credentials against Firebase Authentication in Phase 6 of the plan. "
    "The 'Continue as Guest' option was an intentional design decision: for the academic demo "
    "and presentation, we want to show the map and restaurant features quickly without "
    "needing a real account. The colour palette — a muted steel blue on a light grey "
    "background — was chosen to feel trustworthy and clean, matching the app's purpose "
    "of helping users make decisions about places to eat.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Map & Search Screen
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
set_bg(s8, LIGHT)
add_title_block(s8, "Screen 2 — Map & Search",
                "Live Google Map centred on GPS location with restaurant markers and search")

img8 = os.path.join(SCREENSHOT_DIR, "search feature.png")
add_image(s8, img8, Inches(0.3), Inches(1.5), Inches(2.5), Inches(5.6))

bullet_block(s8, [
    "Google Map loads and centres on the device's live GPS coordinates",
    "FusedLocationProviderClient provides battery-efficient location updates",
    "Location permission requested at runtime via Accompanist Permissions",
    "5 restaurant markers placed at seeded coordinate offsets from user position",
    "Each marker stores a Restaurant object as its tag",
    "",
    "Search bar (top overlay):",
    "  Filters visible markers by name or cuisine type in real-time",
    "  Text input triggers recomposition — only matching markers shown",
    "",
    "Tap any marker  →  navigate to Restaurant Detail screen",
    "Back-stack managed automatically by Jetpack Navigation Compose",
], Inches(3.1), Inches(1.6), Inches(6.6), Inches(5.4),
   font_size=13, title="What This Screen Does")

add_note(s8,
    "This is the core screen of the app. The Google Map is rendered using maps-compose, "
    "which is the official Jetpack Compose wrapper for the Maps SDK. "
    "When the screen first loads it requests ACCESS_FINE_LOCATION permission — if granted, "
    "FusedLocationProviderClient returns the device's GPS position and the camera animates "
    "to that location. The 5 restaurant markers are currently seeded at fixed offsets "
    "from that GPS position, which means they move relative to wherever the user is. "
    "The search bar overlays the map and filters the list of restaurants client-side. "
    "In the next phase, we replace seeded data with a live Places API Nearby Search call "
    "that returns real restaurants within 1500 metres of the user's location.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Restaurant Detail Screen
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
set_bg(s9, LIGHT)
add_title_block(s9, "Screen 3 — Restaurant Detail",
                "Full info, video preview, and navigation — the app's key value screen")

img9 = os.path.join(SCREENSHOT_DIR, "navigation to video and direction.png")
add_image(s9, img9, Inches(0.3), Inches(1.5), Inches(2.5), Inches(5.6))

bullet_block(s9, [
    "TopAppBar with back arrow returns user to Map screen",
    "Restaurant card shows:  cuisine type, address, star rating",
    "Distance row:  pin icon + '450m away' (calculated from GPS delta)",
    "",
    "Watch Video Preview (filled button):",
    "  Opens a YouTube URL via Android Intent → YouTube app or browser",
    "  URL is stored per restaurant in the seeded data map",
    "  No WebView needed — native YouTube app handles playback",
    "",
    "Get Directions (outlined button):",
    "  Builds a maps.google.com URL with the restaurant's lat/lng as destination",
    "  Fires an implicit Intent — Google Maps app opens with turn-by-turn route",
    "",
    "Both buttons demonstrated working on the emulator",
], Inches(3.1), Inches(1.6), Inches(6.6), Inches(5.4),
   font_size=13, title="What This Screen Does")

add_note(s9,
    "The Restaurant Detail screen is where the two core user actions happen: "
    "watch a video to decide if this place looks good, or launch navigation to go there. "
    "The Watch Video button uses a simple Android Intent with ACTION_VIEW and a YouTube URL. "
    "Android's intent resolution system opens the YouTube app if installed, or falls back to "
    "the browser — no custom video player needed, and no YouTube API key required. "
    "The Get Directions button constructs a Google Maps URL with the destination lat/lng "
    "pre-filled. Again, an implicit Intent opens Google Maps in navigation mode. "
    "This pattern — using intents to delegate to installed apps — is efficient, reliable, "
    "and avoids re-implementing functionality that system apps already provide.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — What's Next / Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
set_bg(s10, LIGHT)
add_title_block(s10, "Progress Summary & What's Next",
                "Phases completed vs. remaining work")

phases = [
    ("Phase 1", "Project scaffold, Gradle, API keys",       True),
    ("Phase 2", "Google Maps + GPS centering",              True),
    ("Phase 3", "Restaurant markers (seeded data)",         True),
    ("Phase 4", "Detail screen: video + directions",        True),
    ("Phase 5", "Live Google Places API integration",       False),
    ("Phase 6", "Firebase Auth (login / register)",         False),
    ("Phase 7", "UI polish, loading states, error handling",False),
    ("Phase 8", "Bottom nav bar, Home & Utility screens",   False),
    ("Phase 9", "Final dissertation report writing",        False),
]

col_w = Inches(3.1)
for i, (phase, desc, done) in enumerate(phases):
    row = i // 3
    col = i % 3
    left = Inches(0.3) + col * col_w
    top  = Inches(1.9) + row * Inches(1.55)

    box = s10.shapes.add_shape(1, left, top, Inches(2.9), Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9) if done else RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = GREEN if done else RGBColor(0xCC, 0xCC, 0xCC)

    tf = box.text_frame; tf.word_wrap = True
    p1  = tf.paragraphs[0]
    r1  = p1.add_run()
    r1.text = f"{'DONE' if done else 'TODO'}  {phase}"
    r1.font.size  = Pt(11)
    r1.font.bold  = True
    r1.font.color.rgb = GREEN if done else RGBColor(0x88, 0x88, 0x88)

    p2  = tf.add_paragraph()
    r2  = p2.add_run()
    r2.text = desc
    r2.font.size  = Pt(10)
    r2.font.color.rgb = DARK

add_textbox(s10,
    "Immediate next step:  replace seeded restaurant data with live Google Places API Nearby Search calls.",
    Inches(0.3), Inches(6.55), Inches(9.4), Inches(0.55),
    font_size=13, bold=True, color=STEEL)

add_note(s10,
    "To summarise: Phases 1 through 4 are complete and demonstrated on a running emulator. "
    "The app has a working login screen, a live Google Map centred on the device GPS, "
    "5 restaurant markers that filter by search, and a detail screen where both Watch Video "
    "and Get Directions are functional. "
    "The immediate next step is Phase 5 — replacing the 5 hardcoded restaurants with a "
    "real Places API Nearby Search call, so the app returns actual restaurants near the user. "
    "Phase 6 adds Firebase Authentication to the login screen. "
    "Phase 7 and 8 cover UI polish — adding a bottom navigation bar with Home, Map, Privacy, "
    "and Contact Us tabs — and proper loading/error states. "
    "Phase 9 is the dissertation report, which will document all technical decisions made "
    "during development. Overall the project is on schedule and the core technical risk "
    "— getting a live Google Map running on a real emulator — has been proven out.")

# ── save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "Spot_To_Go_Progress.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
