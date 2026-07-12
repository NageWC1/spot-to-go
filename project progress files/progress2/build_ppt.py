"""
Generates Spot_To_Go_Progress_Report.pptx using the UEA template.
Keeps the exact theme/colours/fonts; only replaces text and adds screenshots.

To regenerate the PPT after edits, run:
    python build_ppt.py
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(os.path.abspath(__file__))
TMPL   = os.path.join(BASE, 'Progress report_UEA_template.pptx')
UI     = os.path.join(BASE, 'ui')
OUTPUT = os.path.join(BASE, 'Spot_To_Go_Progress_Report.pptx')

prs = Presentation(TMPL)

# ── Delete duplicate example slides 8-13 (keep only slides 1-7) ──────────────
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
id_lst = prs.slides._sldIdLst
for _ in range(len(prs.slides) - 7):          # remove from the end
    elem = id_lst[-1]
    r_id = elem.get(f'{{{NS_R}}}id')
    id_lst.remove(elem)
    try:
        prs.slides.part._rels.pop(r_id, None)
    except Exception:
        pass

# ── Helpers ───────────────────────────────────────────────────────────────────
def clear(placeholder):
    """Wipe all runs/breaks from every paragraph, leave one empty paragraph."""
    tc = placeholder.text_frame._txBody
    paras = tc.findall(qn('a:p'))
    for p in paras[1:]:
        tc.remove(p)
    first = paras[0]
    for child in list(first):
        first.remove(child)

def para(tf, text, level=0, bold=False, pt=None, first_para=False):
    """Append (or reuse the first) paragraph in text-frame tf."""
    if first_para:
        p = tf.paragraphs[0]
        p.text = text
        p.level = level
    else:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
    if p.runs:
        run = p.runs[0]
        run.font.bold = bold
        if pt:
            run.font.size = Pt(pt)
    return p

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def linebreak_rect(shape, line1, line2):
    """Put two lines (with a soft line-break) into a Rectangle's text frame."""
    tf = shape.text_frame
    clear(shape)
    p_elem = tf.paragraphs[0]._p
    for r in p_elem.findall(qn('a:r')):
        p_elem.remove(r)
    r1 = etree.SubElement(p_elem, qn('a:r'))
    etree.SubElement(r1, qn('a:t')).text = line1
    etree.SubElement(p_elem, qn('a:br'))
    r2 = etree.SubElement(p_elem, qn('a:r'))
    etree.SubElement(r2, qn('a:t')).text = line2

def pic(slide, fname, left, top, w, h):
    slide.shapes.add_picture(os.path.join(UI, fname),
                             Inches(left), Inches(top),
                             Inches(w), Inches(h))

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides[0]
for sh in s.shapes:
    if sh.name == 'Rectangle 2':
        linebreak_rect(sh, 'Spot To Go  —  Android App',
                           'Progress Report  |  12 July 2026')
    if sh.name == 'Rectangle 3':
        linebreak_rect(sh, 'Nagenthiran Nagarajah',
                           'University of East Anglia  |  12-Jul-2026')

notes(s,
    "Welcome. My name is Nagenthiran and this is my weekly progress report for "
    "Spot To Go — an Android app that lets users discover nearby restaurants on "
    "a Google Map, view details, and watch YouTube or TikTok preview videos. "
    "Since the last report, Firebase Authentication has gone from written-but-inactive "
    "code to a fully working, tested login and registration flow, and the map screen "
    "gained proper navigation so users are never stuck.")

# ── SLIDE 2 — Plan of Last Meeting ────────────────────────────────────────────
s = prs.slides[1]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Plan Of The Last Meeting'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            ("Goals agreed at the previous meeting:", 0, True,  13),
            ("Build all 11 UI screens from the hand-drawn paper prototype", 1, False, 12),
            ("Wire every screen together with a complete navigation graph", 1, False, 12),
            ("Apply a custom Deep Orange colour theme and branded launcher icon", 1, False, 12),
            ("Ensure consistent background and font styling across all screens", 1, False, 12),
            ("Commit all code and push to GitHub (NageWC1/spot-to-go)", 1, False, 12),
            ("", 0, False, 10),
            ("Planned for the next phase (discussed at last meeting):", 0, True,  13),
            ("Firebase Authentication — real registration and login", 1, False, 12),
            ("Google Places Nearby Search API — replace seed data with live restaurants", 1, False, 12),
            ("Live search bar connected to the Places API keyword parameter", 1, False, 12),
        ]
        for i, (text, lvl, bold, pt) in enumerate(rows):
            para(tf, text, lvl, bold, pt, first_para=(i == 0))

notes(s,
    "At our last meeting we agreed on five concrete goals for this week. "
    "The main goal was implementing all 11 screens visible in the paper prototype. "
    "We also agreed to apply a branded colour theme and icon so the app looks "
    "polished on a real device. "
    "Looking ahead, we had already discussed the next phase: Firebase for auth "
    "and Google Places for live restaurant data. Those form the basis of next week's tasks.")

# ── SLIDE 3 — Progress Overview ───────────────────────────────────────────────
s = prs.slides[2]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Progress Overview For The Last Week'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        done = [
            "Splash Screen — 2-second auto-navigate, loading indicator",
            "Home Page — search bar, Explore Nearby button, 5-tab bottom navigation",
            "Login Page — email/password, Forgot Password, Guest mode, Register link",
            "Registration Page — full validation, Privacy Policy checkbox, error messages",
            "Map Screen — GPS centring, 5 seed restaurant markers, live search filter",
            "Restaurant Detail — name, rating, distance, cuisine type, action buttons",
            "Video Screen — YouTube placeholder player, Watch / TikTok dual buttons",
            "TikTok Link Page — branded dark preview, like/comment/share, OPEN IN TIKTOK",
            "Directions Screen — Car / Bus / Walk mode chips, live map, START button",
            "Contact Us Page — validated form, Snackbar success confirmation",
            "Privacy Policy — scrollable 5-section document",
            "Deep Orange theme (#FF5722) + gradient launcher icon  →  GitHub commit c87e960",
            "Firebase Auth SDK wired: AuthRepository, real login/register, auth gate, loading spinner  →  commit 1b9c28c",
            "Firebase project connected — google-services.json added, Maps SDK key configured, first successful build",
            "Login mandatory before the map — MapScreen redirects to Login unless a Firebase session is active",
            "Register/Login UX fixes — password show/hide toggle, keyboard no longer hides Password/Confirm fields",
            "Map screen navigation fixed — added bottom nav bar (Home / Map / Contact / Privacy / Logout)",
            "Login state now consistent across screens — Home's nav bar reflects real signed-in/out status",
        ]
        for i, text in enumerate(done):
            para(tf, f'✓  {text}', 0, False, 12, first_para=(i == 0))

notes(s,
    "All goals from the last meeting were completed, and the Firebase Authentication work "
    "went further than planned. "
    "The tick list covers every screen from the paper prototype, the orange theme and icon, "
    "and now a fully activated auth backend. "
    "Once google-services.json was added and the Maps API key configured, the first real "
    "build succeeded end-to-end. "
    "Testing on a physical device surfaced three UX bugs that were fixed this session: "
    "the keyboard was hiding the password fields (fixed with imePadding), there was no way "
    "to verify a typed password before submitting (added a show/hide toggle), and the map "
    "screen was a dead end with no way back to other screens except the phone's back button, "
    "which just exited the app (fixed by adding a bottom navigation bar consistent with Home). "
    "Login is now mandatory before reaching the map, and the two nav bars agree on whether "
    "the user is actually signed in.")

# ── SLIDE 4 — All 11 Screenshots ──────────────────────────────────────────────
s = prs.slides[3]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Progress For The Last Week — UI Screens'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        p = tf.paragraphs[0]
        p.text = 'All 11 screens built in Jetpack Compose, matching the hand-drawn prototype.'
        if p.runs:
            p.runs[0].font.italic = True
            p.runs[0].font.size   = Pt(11)

# Image grid: row 1 (6 images), row 2 (5 images)
ROW1 = ['splash screen.jpeg', 'home page.jpeg', 'login page.jpeg',
        'registration page.jpeg', 'map page.jpeg', 'hotel review page.jpeg']
ROW2 = ['watch video page.jpeg', 'tiktok launch page.jpeg',
        'direction screen.jpeg', 'contact us page.jpeg', 'privacy page.jpeg']

W, H, GAP = 1.0, 2.05, 0.10
TOP1 = 2.2
TOP2 = TOP1 + H + 0.22

def row_left(n):
    return 0.5 + (9.0 - (n * W + (n - 1) * GAP)) / 2

for i, f in enumerate(ROW1):
    pic(s, f, row_left(len(ROW1)) + i * (W + GAP), TOP1, W, H)
for i, f in enumerate(ROW2):
    pic(s, f, row_left(len(ROW2)) + i * (W + GAP), TOP2, W, H)

notes(s,
    "Walk through the two rows of screenshots. "
    "Row 1 (left to right): Splash Screen with the orange gradient and location pin, "
    "Home Page with the search bar and Explore Nearby button, "
    "Login Page, Registration Page with the checkbox and validation, "
    "Map Screen showing live Google Map with restaurant markers, "
    "and the Restaurant Detail page. "
    "Row 2: Video Screen with the YouTube thumbnail and TikTok button, "
    "TikTok Link Page with the branded dark UI, "
    "Directions Screen with the transport mode chips and live map, "
    "Contact Us form, and the Privacy Policy. "
    "Every screen uses Material 3 with our Deep Orange theme. "
    "Navigation is handled entirely by Jetpack Compose Navigation — no Activities, "
    "just composable screens connected in a single NavHost.")

# ── SLIDE 5 — Built-in Backends ───────────────────────────────────────────────
s = prs.slides[4]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Progress For The Last Week — Built-In Backends'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            ("The app has 7 backend systems — 6 live today, 1 ready to activate:", 0, True, 13),
            ("1. GPS  (FusedLocationProviderClient)", 1, True,  12),
            ("Real device coordinates centre the map and position all seed restaurants around the user.", 2, False, 10),
            ("2. Google Maps SDK  —  live map tiles, markers, camera animation", 1, True,  12),
            ("Fully live Google Map. Tap a marker → restaurant name and rating appear instantly.", 2, False, 10),
            ("3. RestaurantRepository  —  in-memory data layer (temporary database)", 1, True,  12),
            ("5 seed restaurants with name, cuisine, rating, distance, YouTube URL. Replaced by Places API next week.", 2, False, 10),
            ("4. Jetpack Navigation  —  back-stack and route management", 1, True,  12),
            ("All 11 screen routes managed centrally. Back button, pop-up stack, and deep links all work.", 2, False, 10),
            ("5. Android Intent System  —  external app launcher", 1, True,  12),
            ("Opens YouTube for videos, TikTok for restaurant search, Google Maps for turn-by-turn navigation.", 2, False, 10),
            ("6. Compose State  (remember / mutableStateOf)", 1, True,  12),
            ("Search queries update markers in real-time. Selected restaurant survives Map → Detail → Video.", 2, False, 10),
            ("7. Firebase AuthRepository  —  live and tested on a physical device", 1, True,  12),
            ("register() / login() / logout() / isLoggedIn wired to Firebase and confirmed working end-to-end. Map is gated behind login.", 2, False, 10),
        ]
        for i, (text, lvl, bold, pt) in enumerate(rows):
            para(tf, text, lvl, bold, pt, first_para=(i == 0))

notes(s,
    "This slide addresses a common question: is there a real backend behind this app? "
    "The answer is yes, all seven systems are now live, including Firebase Authentication "
    "which was activated and tested this week. "
    "\n"
    "GPS: the FusedLocationProviderClient returns the device's real coordinates. "
    "The map truly centres on where you are. "
    "\n"
    "Google Maps SDK: the map tiles are live from Google's servers. "
    "Markers, camera animation, and the my-location button are all real SDK features. "
    "\n"
    "RestaurantRepository acts as a temporary in-memory database. "
    "It will be replaced by the Places API next week. "
    "\n"
    "Jetpack Navigation manages all 11 routes — the back stack is real, "
    "pressing Back from Detail returns to Map correctly. "
    "\n"
    "The Intent system connects to real external apps: tapping Watch Video "
    "opens the actual YouTube video; GET DIRECTIONS opens Google Maps navigation. "
    "\n"
    "Finally, Compose state keeps the search query and selected restaurant "
    "alive across recompositions — the map search filter works live today.")

# ── SLIDE 6 — Actions for Next Week (Milestones) ─────────────────────────────
s = prs.slides[5]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Actions for the Next Week'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            ("Milestones", 0, True,  14),
            ("Dissertation Proposal draft — submitted to supervisor", 1, False, 12),
            ("Progress Report — continue weekly reporting cycle", 1, False, 12),
            ("Backend Phase — Firebase Auth live; Places API now the active target", 1, False, 12),
            ("", 0, False, 10),
            ("Starting Point for Next Session", 0, True,  14),
            ("data/Restaurant.kt → RestaurantRepository.getSeedRestaurants() is where work begins:", 1, False, 12),
            ("replace its hardcoded list with a new PlacesRepository backed by a live Nearby Search call", 2, False, 11),
            ("MapScreen.kt already consumes RestaurantRepository, so the UI needs no rework — only the data source changes", 2, False, 11),
            ("", 0, False, 10),
            ("Tasks for Next Session", 0, True,  14),
            ("Task 1  —  Google Places Nearby Search API (real restaurant data)", 1, False, 12),
            ("Task 2  —  Live Search connected to Places API keyword parameter", 1, False, 12),
            ("Task 3  —  Double-back-to-exit confirmation on the map screen (UX polish)", 1, False, 12),
        ]
        for i, (text, lvl, bold, pt) in enumerate(rows):
            para(tf, text, lvl, bold, pt, first_para=(i == 0))

notes(s,
    "Three milestones are active. "
    "The dissertation proposal draft has been submitted. "
    "Weekly reporting continues. "
    "The backend phase has moved on from auth — that is now live — to the Places API. "
    "\n"
    "To be concrete about where next session starts: data/Restaurant.kt currently defines "
    "RestaurantRepository.getSeedRestaurants(), which returns five hardcoded restaurants "
    "offset from the user's GPS position. That function is the exact starting point — it "
    "gets replaced by a new PlacesRepository that calls the real Places Nearby Search API. "
    "Because MapScreen.kt already reads from RestaurantRepository rather than holding data "
    "itself, no UI changes are needed; only the data source underneath it changes. "
    "\n"
    "Task 1 is now the immediate next action: Firebase Auth activation is complete and tested, "
    "so the next backend milestone is replacing the five hardcoded seed restaurants with a real "
    "Google Places Nearby Search call. "
    "\n"
    "Task 2 follows once Places API data is flowing: the live search bar sends the typed keyword "
    "to the API with a debounce instead of filtering the seed list client-side. "
    "\n"
    "Task 3 is a small UX polish item raised during device testing: pressing back on the map "
    "currently exits the app immediately since it is the root screen after login — a "
    "press-back-again-to-exit confirmation would be friendlier than an instant close.")

# ── SLIDE 7 — Actions for Next Week (Detail) ─────────────────────────────────
s = prs.slides[6]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Actions for the Next Week — Task Detail'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            ("Firebase Auth — COMPLETE  (activated and tested this session)", 0, True,  13),
            ("✓  google-services.json added to app/ and Maps SDK API key configured", 1, False, 11),
            ("✓  Email/Password sign-in enabled in Firebase Console", 1, False, 11),
            ("✓  First successful build — ./gradlew assembleDebug — confirmed on device", 1, False, 11),
            ("✓  Register → auto-login → Map flow tested end-to-end on a physical phone", 1, False, 11),
            ("✓  Fixed: keyboard hid Password/Confirm fields — added imePadding on Login/Register", 1, False, 11),
            ("✓  Fixed: no way to verify typed password — added show/hide toggle", 1, False, 11),
            ("✓  Fixed: map screen had no way back — added bottom nav (Home/Map/Contact/Privacy/Logout)", 1, False, 11),
            ("✓  Fixed: Home nav bar showed 'Login' even when signed in — now reflects real auth state", 1, False, 11),
            ("Task 1 — Google Places Nearby Search API  (IMMEDIATE NEXT STEP)", 0, True,  13),
            ("Add Places SDK to libs.versions.toml and app/build.gradle.kts", 1, False, 11),
            ("Create PlacesRepository calling Nearby Search (type: restaurant, radius: 1500 m)", 1, False, 11),
            ("Replace RestaurantRepository seed data with live API results on MapScreen", 1, False, 11),
            ("Task 2 — Live Search Functionality", 0, True,  13),
            ("Pass MapScreen search bar text as keyword param in Places API call", 1, False, 11),
            ("Refresh markers in real-time as the user types (debounce 400 ms)", 1, False, 11),
            ("Show CircularProgressIndicator while API request is in progress", 1, False, 11),
            ("Task 3 — Map Back-Press UX Polish", 0, True,  13),
            ("Add press-back-again-to-exit confirmation on the map screen", 1, False, 11),
        ]
        for i, (text, lvl, bold, pt) in enumerate(rows):
            para(tf, text, lvl, bold, pt, first_para=(i == 0))

notes(s,
    "Firebase Authentication is now fully complete, not just code-written. "
    "google-services.json was added, the Maps API key was configured, and Email/Password "
    "sign-in was enabled in the Firebase Console. The first build succeeded, and the full "
    "register-to-map flow was tested on a physical device rather than just an emulator. "
    "That real-device testing is what surfaced four UX bugs, all now fixed: the keyboard "
    "hiding the password fields, no way to double-check a typed password, the map screen "
    "being a dead end, and the Home nav bar not reflecting the real sign-in state. "
    "\n"
    "Task 1 — Places API: "
    "Now that auth is confirmed working, we add the Places SDK, "
    "create PlacesRepository hitting Nearby Search within 1500 metres, "
    "and replace the five hardcoded seed restaurants with live data. "
    "\n"
    "Task 2 — Live Search: "
    "The search bar text becomes the keyword parameter in the Places call, "
    "with a 400ms debounce and a loading indicator while the request is in flight. "
    "\n"
    "Task 3 — Back-press polish: "
    "A small but user-visible fix identified during testing — pressing back on the map "
    "currently exits the app immediately since it has no parent screen in the back stack.")

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'Done — {len(prs.slides)} slides saved to:')
print(OUTPUT)
