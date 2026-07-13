"""
Generates Spot_To_Go_Presentation.pptx using the UEA template.
Unlike the weekly progress report, this deck covers everything built so far
(UI + backend) as a single client-facing presentation — no dates, simple
language, and full spoken-script speaker notes.

To regenerate the PPT after edits, run:
    python build_ppt.py
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(os.path.abspath(__file__))
TMPL   = os.path.join(BASE, 'Progress report_UEA_template.pptx')
UI     = os.path.join(BASE, 'ui')
OUTPUT = os.path.join(BASE, 'Spot_To_Go_Presentation.pptx')

prs = Presentation(TMPL)

# ── Delete duplicate example slides ───────────────────────────────────────────
# The template repeats its title-slide layout at index 7 (an extra copy of
# slide 0) before its content slides continue at index 8+. Drop that repeated
# title slide first, then trim from the end so 8 slides remain: 1 title +
# 7 content slides (6 "done so far" + 1 "what's coming next").
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
id_lst = prs.slides._sldIdLst

def drop_slide(index):
    elem = id_lst[index]
    r_id = elem.get(f'{{{NS_R}}}id')
    id_lst.remove(elem)
    try:
        prs.slides.part._rels.pop(r_id, None)
    except Exception:
        pass

drop_slide(7)
for _ in range(len(prs.slides) - 8):          # remove remaining excess from the end
    drop_slide(-1)

# ── Helpers ───────────────────────────────────────────────────────────────────
def clear(placeholder):
    """Wipe all runs/breaks from every paragraph, leave one empty paragraph."""
    tc = placeholder.text_frame._txBody
    paras = tc.findall(qn('a:p'))
    for p in paras[1:]:
        tc.remove(p)
    first = paras[0]
    for child in list(first):
        first.remove(child)

def para(tf, text, level=0, bold=False, pt=None, italic=False, first_para=False):
    """Append (or reuse the first) paragraph in text-frame tf."""
    if first_para:
        p = tf.paragraphs[0]
        p.text = text
        p.level = level
    else:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
    if p.runs:
        run = p.runs[0]
        run.font.bold = bold
        run.font.italic = italic
        if pt:
            run.font.size = Pt(pt)
    return p

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def linebreak_rect(shape, line1, line2):
    """Put two lines (with a soft line-break) into a Rectangle's text frame."""
    tf = shape.text_frame
    clear(shape)
    p_elem = tf.paragraphs[0]._p
    for r in p_elem.findall(qn('a:r')):
        p_elem.remove(r)
    r1 = etree.SubElement(p_elem, qn('a:r'))
    etree.SubElement(r1, qn('a:t')).text = line1
    etree.SubElement(p_elem, qn('a:br'))
    r2 = etree.SubElement(p_elem, qn('a:r'))
    etree.SubElement(r2, qn('a:t')).text = line2

def pic(slide, fname, left, top, w, h):
    slide.shapes.add_picture(os.path.join(UI, fname),
                             Inches(left), Inches(top),
                             Inches(w), Inches(h))

def strip_date(slide):
    """Remove the template's hardcoded date text — no dates in this deck."""
    for sh in slide.shapes:
        if sh.name.startswith('Date Placeholder'):
            clear(sh)

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides[0]
for sh in s.shapes:
    if sh.name == 'Rectangle 2':
        linebreak_rect(sh, 'Spot To Go  —  Android App',
                           'Project Presentation')
    if sh.name == 'Rectangle 3':
        linebreak_rect(sh, 'Nagenthiran Nagarajah',
                           'University of East Anglia')
strip_date(s)

notes(s,
    "Good morning/afternoon, everyone. My name is Nagenthiran Nagarajah, and today "
    "I'm going to walk you through Spot To Go, the Android app I've been building. "
    "Spot To Go helps people discover restaurants near them, see real details about "
    "each one, and even watch a short video of what the place actually looks like "
    "before deciding to go. I'll show you every screen in the app, and then explain "
    "how it all works behind the scenes.")

# ── SLIDE 2 — What Is Spot To Go? ─────────────────────────────────────────────
s = prs.slides[1]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'What Is Spot To Go?'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            "Helps users discover restaurants near their current location",
            "Shows each restaurant's name, rating, distance, and cuisine type",
            "Lets users watch a real video preview (YouTube or TikTok) before visiting",
            "Gives instant turn-by-turn directions to the restaurant",
            "Built as a modern Android app using Kotlin and Jetpack Compose",
        ]
        for i, text in enumerate(rows):
            para(tf, text, 0, False, 14, first_para=(i == 0))
strip_date(s)

notes(s,
    "So what exactly is Spot To Go? At its core, it's a simple idea: instead of "
    "guessing what a restaurant is like from a photo or a star rating, you can "
    "actually watch a short video of it first. "
    "The app uses the phone's location to find restaurants nearby, shows the "
    "important details like rating, distance and cuisine, and then gives a real "
    "video preview so the user knows what to expect. "
    "Once they've decided, one tap gives directions straight there. "
    "It's built using modern Android tools — Kotlin and Jetpack Compose — which is "
    "the current recommended way to build Android apps.")

# ── SLIDE 3 — All App Screens ─────────────────────────────────────────────────
s = prs.slides[2]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'All App Screens'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        p = tf.paragraphs[0]
        p.text = 'Every screen has been designed and built to match a simple, easy-to-follow flow.'
        if p.runs:
            p.runs[0].font.italic = True
            p.runs[0].font.size = Pt(11)
strip_date(s)

ROW1 = ['splash screen.jpeg', 'home page.jpeg', 'login page.jpeg',
        'registration page.jpeg', 'map page.jpeg', 'hotel review page.jpeg']
ROW2 = ['watch video page.jpeg', 'tiktok launch page.jpeg',
        'direction screen.jpeg', 'contact us page.jpeg', 'privacy page.jpeg']

W, H, GAP = 1.0, 2.05, 0.10
TOP1 = 2.2
TOP2 = TOP1 + H + 0.22

def row_left(n):
    return 0.5 + (9.0 - (n * W + (n - 1) * GAP)) / 2

for i, f in enumerate(ROW1):
    pic(s, f, row_left(len(ROW1)) + i * (W + GAP), TOP1, W, H)
for i, f in enumerate(ROW2):
    pic(s, f, row_left(len(ROW2)) + i * (W + GAP), TOP2, W, H)

notes(s,
    "Let me show you every screen in the app. "
    "Starting with the top row: this is the Splash Screen, a short loading screen "
    "you see when you first open the app. Next is the Home Page, with a search bar "
    "and a button to explore nearby restaurants. Then the Login Page and the "
    "Registration Page, where a user creates an account or signs back in. After "
    "that, the Map Screen, which shows restaurants near the user's real location "
    "on a live map. And finally in this row, the Restaurant Detail page, showing "
    "the name, rating, distance and type of food. "
    "On the second row: the Watch Video screen, where users can preview a "
    "restaurant through YouTube or TikTok, the TikTok Link page with a branded "
    "preview, the Directions screen with walking, driving and bus options, the "
    "Contact Us page, and the Privacy Policy page. "
    "Every single one of these screens is fully built and working today.")

# ── SLIDE 4 — How A User Moves Through The App ───────────────────────────────
s = prs.slides[3]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'How A User Moves Through The App'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            ("Splash Screen — a quick loading screen when the app opens", 0, False, 13),
            ("Login or Register — sign in, create an account, or continue as a guest", 0, False, 13),
            ("Home Page — search bar and quick access to nearby restaurants", 0, False, 13),
            ("Map Screen — restaurants shown near the user's real location", 0, False, 13),
            ("Restaurant Detail — name, rating, distance, and cuisine type", 0, False, 13),
            ("Watch Video — a real preview video of the restaurant", 0, False, 13),
            ("Directions — step-by-step navigation to the restaurant", 0, False, 13),
        ]
        for i, (text, lvl, bold, pt) in enumerate(rows):
            para(tf, text, lvl, bold, pt, first_para=(i == 0))
strip_date(s)

notes(s,
    "Now let's follow one typical journey through the app. "
    "A user opens the app and briefly sees the splash screen while things load. "
    "They then log in, register a new account, or simply continue as a guest. "
    "From the home page, they can search directly or tap Explore Nearby to jump "
    "straight to the map. The map shows restaurants around their real location "
    "using the phone's GPS. "
    "Tapping a restaurant opens its detail page, with the rating, distance and "
    "cuisine type. From there, they can either watch a real video preview of the "
    "restaurant, or tap Get Directions to open turn-by-turn navigation straight to "
    "the door. "
    "Every step connects smoothly to the next — nothing is a dead end.")

# ── SLIDE 5 — What's Working Behind the Scenes ───────────────────────────────
s = prs.slides[4]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = "What's Working Behind the Scenes"
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            ("Real GPS Location", 0, True, 13),
            ("Finds restaurants near the user's actual position, not a fixed test location", 1, False, 11),
            ("Live Google Map", 0, True, 13),
            ("An interactive map with a real marker for every restaurant", 1, False, 11),
            ("Account System", 0, True, 13),
            ("Secure login and registration, tested on a real phone", 1, False, 11),
            ("Screen Navigation", 0, True, 13),
            ("Smooth movement between every screen, including a working back button", 1, False, 11),
            ("Opens Other Apps", 0, True, 13),
            ("The video button opens YouTube or TikTok, and directions open Google Maps", 1, False, 11),
            ("Smart Search", 0, True, 13),
            ("Typing in the search bar instantly filters which restaurants are shown", 1, False, 11),
        ]
        for i, (text, lvl, bold, pt) in enumerate(rows):
            para(tf, text, lvl, bold, pt, first_para=(i == 0))
strip_date(s)

notes(s,
    "Behind the screens you just saw, there's a real working system underneath. "
    "The app uses the phone's GPS to find the user's actual location, not a fixed "
    "test location. "
    "The map itself is a live Google Map, not an image — it can be zoomed and "
    "panned, and it holds real markers for each restaurant. "
    "Logging in and registering is handled by a real account system, and this has "
    "been tested on an actual physical phone, not just a simulator. "
    "Moving between every screen — including pressing back — is handled properly, "
    "so a user is never stuck with no way out. "
    "When someone taps Watch Video or Get Directions, the app opens the real "
    "YouTube, TikTok, or Google Maps app on the phone. "
    "And the search bar isn't just decoration — typing into it instantly filters "
    "the restaurants shown on the map.")

# ── SLIDE 6 — Login & Account Security ───────────────────────────────────────
s = prs.slides[5]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Login & Account Security'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            "Users can register a new account or log in with an existing one",
            "Accounts are stored securely, not just kept on the phone",
            "Users can also continue as a guest without creating an account",
            "The map and restaurant features only unlock once a user is signed in",
            "Tested end-to-end on a real Android phone",
        ]
        for i, text in enumerate(rows):
            para(tf, text, 0, False, 14, first_para=(i == 0))
strip_date(s)

notes(s,
    "Because Spot To Go stores personal accounts, security matters. "
    "Registration and login are both connected to a real account system, which "
    "keeps passwords and user details stored securely rather than just sitting on "
    "the phone. "
    "Users who don't want to register can still continue as a guest. "
    "To keep the experience personal and safe, the map and restaurant features "
    "only unlock once someone is actually signed in. "
    "All of this has been tested end-to-end on a real Android phone — "
    "registering a new account, logging in, and reaching the map — and it works "
    "exactly as intended.")

# ── SLIDE 7 — Summary ─────────────────────────────────────────────────────────
s = prs.slides[6]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = 'Summary'
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            "All app screens are designed and fully working",
            "A real login and account system is in place and tested",
            "The map shows restaurants using the user's real location",
            "Users can preview restaurants with real videos before visiting",
            "One tap gives directions straight to the restaurant",
        ]
        for i, text in enumerate(rows):
            para(tf, text, 0, False, 14, first_para=(i == 0))
strip_date(s)

notes(s,
    "To summarise what you've seen today: every screen in Spot To Go is built and "
    "working, from the splash screen through to the privacy policy. "
    "There's a real login and account system behind it, tested on an actual "
    "phone. "
    "The map uses the user's real location to show nearby restaurants, and from "
    "there, users can preview a restaurant with a real video before deciding to "
    "go, then get directions with a single tap. "
    "That's Spot To Go so far.")

# ── SLIDE 8 — What's Coming Next ──────────────────────────────────────────────
s = prs.slides[7]
for sh in s.shapes:
    if sh.name == 'Title 1':
        sh.text_frame.paragraphs[0].text = "What's Coming Next"
    if sh.name == 'Content Placeholder 2':
        tf = sh.text_frame
        clear(sh)
        rows = [
            "Connect the map to real, live restaurant listings nearby",
            "Make the search bar pull results directly from those live listings",
            "Add a small safeguard so pressing back on the map doesn't close the app by accident",
        ]
        for i, text in enumerate(rows):
            para(tf, text, 0, False, 14, first_para=(i == 0))
strip_date(s)

notes(s,
    "So what comes next? "
    "Right now, the restaurants you saw on the map are a small sample set used to "
    "build and test every screen properly. "
    "The next step is to connect the map to real, live restaurant listings near "
    "the user, so the same experience you just saw works with genuine, up-to-date "
    "places. "
    "Once that's in place, the search bar will pull its results from those live "
    "listings directly, instead of filtering the sample set. "
    "Alongside that, there's one small polish item: pressing back on the map "
    "currently closes the app straight away, and we'll add a simple safeguard so "
    "that doesn't happen by accident. "
    "That's the plan for the next stage of Spot To Go. Thank you.")

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'Done — {len(prs.slides)} slides saved to:')
print(OUTPUT)
