from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

IMG_DIR = r"C:\Users\Akaam Zain\Desktop\NageWorks\spot-to-go\final report\images"
OUT = r"C:\Users\Akaam Zain\Desktop\NageWorks\spot-to-go\project progress files\progress3\Spot_To_Go_Final_Presentation.pptx"

ORANGE = RGBColor(0xFF, 0x57, 0x22)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.0,
             font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, size, color, spacing=14,
                 bullet_color=ORANGE, numbered=False, bold_lead=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        # bullet marker run
        marker = f"{i+1}.  " if numbered else "\u25CF   "
        r1 = p.add_run()
        r1.text = marker
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r1.font.name = FONT
        if numbered:
            r1.font.size = Pt(size)
        else:
            r1.font.size = Pt(size - 4)
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.bold = False
        r2.font.color.rgb = color
        r2.font.name = FONT
    return tb


def add_kicker_and_title(slide, kicker, title, num):
    # top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.12), ORANGE)
    add_text(slide, Inches(0.7), Inches(0.35), Inches(6), Inches(0.4),
              kicker.upper(), 14, ORANGE, bold=True)
    add_text(slide, Inches(0.7), Inches(0.68), Inches(10.5), Inches(0.9),
              title, 32, DARK, bold=True)
    # slide number bottom right
    add_text(slide, SW - Inches(1.0), SH - Inches(0.55), Inches(0.6), Inches(0.4),
              str(num), 12, GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.7), SH - Inches(0.55), Inches(2), Inches(0.4),
              "Spot To Go", 12, GRAY)


def content_slide(kicker, title, bullets, num, numbered=False):
    s = add_slide()
    set_bg(s, WHITE)
    add_kicker_and_title(s, kicker, title, num)
    add_bullets(s, Inches(0.9), Inches(2.05), Inches(11.2), Inches(4.6), bullets,
                22, DARK, spacing=18, numbered=numbered)
    return s


def screens_slide(kicker, title, filenames, captions, note, num):
    s = add_slide()
    set_bg(s, WHITE)
    add_kicker_and_title(s, kicker, title, num)

    n = len(filenames)
    img_h = Inches(3.7)
    gap = Inches(0.55)
    # compute width per image keeping aspect ratio, then total width to center
    from PIL import Image
    widths = []
    for f in filenames:
        im = Image.open(f"{IMG_DIR}\\{f}")
        w, h = im.size
        widths.append(img_h * (w / h))
    total_w = sum(widths, Emu(0)) + gap * (n - 1)
    start_x = (SW - total_w) / 2
    top = Inches(1.85)
    x = start_x
    for f, cap, w in zip(filenames, captions, widths):
        s.shapes.add_picture(f"{IMG_DIR}\\{f}", x, top, height=img_h, width=w)
        add_text(s, x - Inches(0.15), top + img_h + Inches(0.10), w + Inches(0.3), Inches(0.4),
                  cap, 15, DARK, bold=True, align=PP_ALIGN.CENTER)
        x += w + gap

    add_text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.45),
              note, 16, GRAY, italic=True, align=PP_ALIGN.CENTER)
    return s


def cover_slide():
    s = add_slide()
    set_bg(s, ORANGE)
    add_text(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.3),
              "Spot To Go", 60, WHITE, bold=True)
    add_text(s, Inches(1.0), Inches(3.55), Inches(11.3), Inches(0.7),
              "Discover Restaurants. Watch Real Videos. Get There Fast.", 22, WHITE)
    add_rect(s, Inches(1.0), Inches(4.5), Inches(2.2), Pt(3), WHITE)
    add_text(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(0.5),
              "Final Project Presentation", 18, WHITE, bold=True)
    add_text(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.5),
              "University of East Anglia  \u2022  August 2026", 16, WHITE)
    return s


def closing_slide():
    s = add_slide()
    set_bg(s, ORANGE)
    add_text(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.2),
              "Thank You", 56, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(3.9), Inches(11.3), Inches(0.7),
              "Questions?", 26, WHITE, align=PP_ALIGN.CENTER)
    return s


# ---------------- Build deck ----------------

cover_slide()  # 1

content_slide(
    "The Problem", "Choosing Where To Eat Is Still A Guess",
    [
        "Star ratings don't show what a place actually feels like",
        "Short text reviews are easy to skim past \u2014 or fake",
        "You still don't know if it's \u201cyour kind of place\u201d until you're there",
    ], 2)

content_slide(
    "The Idea", "Spot To Go",
    [
        "Shows nearby restaurants on a live map",
        "Tap one to watch a real video preview \u2014 see it before you go",
        "One more tap gives you directions there",
    ], 3)

content_slide(
    "What The App Does", "Five Things It Lets You Do",
    [
        "Sign in securely",
        "See restaurants near you on a live map",
        "Search by name or keyword",
        "Watch a real video preview",
        "Get instant directions",
    ], 4, numbered=True)

screens_slide(
    "The App", "Getting In",
    ["splash screen.jpeg", "registration page.jpeg", "login page.jpeg"],
    ["Splash", "Register", "Login"],
    "New user? Create an account first.  Already have one? Sign straight in.",
    5)

screens_slide(
    "The App", "Finding A Restaurant",
    ["home page.jpeg", "map page.jpeg", "hotel review page.jpeg"],
    ["Home", "Map", "Details"],
    "Search or explore nearby, tap a marker, see the full details.",
    6)

screens_slide(
    "The App", "Deciding & Going",
    ["watch video page.jpeg", "tiktok launch page.jpeg", "direction screen.jpeg"],
    ["Video", "TikTok", "Directions"],
    "Watch a real preview, then get turn-by-turn directions.",
    7)

content_slide(
    "Under The Hood", "In Plain Terms",
    [
        "Built for Android using Kotlin",
        "Uses the phone's real GPS location",
        "Real Google Map with live restaurant markers",
        "Secure sign-in \u2014 the app never stores your password itself",
        "Video and Directions buttons open the real YouTube, TikTok and Google Maps apps",
    ], 8)

content_slide(
    "Testing", "Tested On A Real Phone",
    [
        "Every screen tested end-to-end on a physical Android device",
        "5 real issues found \u2014 and fixed \u2014 during testing",
        "Example: the keyboard was covering the password box \u2014 fixed",
        "Example: no way back from the map screen \u2014 added a menu bar",
        "Every screen connects correctly to the next",
    ], 9)

content_slide(
    "What's Next", "What's Left To Do",
    [
        "Connect to real, live nearby restaurants (Google Places API)",
        "Make search pull results from those live listings",
        "A few small polish fixes, plus more testing",
    ], 10, numbered=True)

content_slide(
    "Summary", "Where Things Stand",
    [
        "All 11 screens built and connected",
        "Real login, tested on a real phone",
        "Map uses the user's real GPS location",
        "Real video previews and real directions, one tap away",
        "One clear next step: live restaurant data",
    ], 11)

closing_slide()  # 12

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
